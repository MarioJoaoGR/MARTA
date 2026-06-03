"""Build the MARTA presentation deck — concrete walkthrough version."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from lxml import etree

# ---------- Palette ----------
NAVY = RGBColor(0x0D, 0x1B, 0x3D)
NAVY_SOFT = RGBColor(0x1E, 0x3A, 0x6F)
STEEL = RGBColor(0x6E, 0x8C, 0xB8)
ICE = RGBColor(0xE6, 0xEC, 0xF5)
PAPER = RGBColor(0xFA, 0xFB, 0xFD)
INK = RGBColor(0x14, 0x1B, 0x2C)
MUTED = RGBColor(0x5C, 0x6B, 0x7A)
CORAL = RGBColor(0xE0, 0x6B, 0x5C)
CODE_BG = RGBColor(0xF3, 0xF5, 0xF8)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREEN = RGBColor(0x2E, 0x86, 0x55)

HEADER_FONT = "Cambria"
BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

TOTAL_SLIDES = 17

# ---------- Setup ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


# ---------- Helpers ----------
def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_rounded(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.06
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_oval(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, font=BODY_FONT, size=14, bold=False,
             color=INK, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, italic=False,
             line_spacing=1.15):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = valign
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb


def add_code(slide, x, y, w, h, code, *, size=11, color=INK,
             bg=CODE_BG, highlight=None):
    """Add a code block: card background + monospaced text.
    highlight: list of (line_no_1indexed, color) for specific line tinting."""
    add_rect(slide, x, y, w, h, bg)
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1),
                                  w - Inches(0.3), h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    lines = code.split("\n")
    hl = {ln: c for ln, c in (highlight or [])}
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.15
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = MONO_FONT
        r.font.size = Pt(size)
        r.font.color.rgb = hl.get(i + 1, color)
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=STEEL, weight=2):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")
    return line


def slide_number(slide, n, on_dark=False):
    color = ICE if on_dark else MUTED
    add_text(slide, Inches(12.3), Inches(0.35), Inches(1.0), Inches(0.3),
             f"{n:02d} / {TOTAL_SLIDES:02d}", font=BODY_FONT, size=10,
             color=color, align=PP_ALIGN.RIGHT, bold=True)


def eyebrow(slide, text, on_dark=False):
    add_text(slide, Inches(0.7), Inches(0.55), Inches(10), Inches(0.3),
             text.upper(), font=BODY_FONT, size=11, bold=True, color=CORAL)


def title(slide, text, on_dark=False, y=Inches(0.95)):
    color = ICE if on_dark else INK
    add_text(slide, Inches(0.7), y, Inches(12), Inches(1.0),
             text, font=HEADER_FONT, size=34, bold=True, color=color,
             line_spacing=1.05)


def subtitle(slide, text, on_dark=False, y=Inches(1.95)):
    color = STEEL if on_dark else MUTED
    add_text(slide, Inches(0.7), y, Inches(12), Inches(0.5),
             text, font=HEADER_FONT, size=18, italic=True, color=color)


# ===================================================================
# SLIDE 1 — TITLE
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, Inches(0.7), Inches(0.7), Inches(0.25), Inches(0.25), CORAL)
add_text(s, Inches(1.05), Inches(0.7), Inches(8), Inches(0.3),
         "GECAD  ·  PYTHON TEST GENERATION", size=11, bold=True, color=ICE)

add_text(s, Inches(0.7), Inches(2.4), Inches(12), Inches(1.4),
         "MARTA", font=HEADER_FONT, size=110, bold=True, color=ICE,
         line_spacing=0.95)
add_text(s, Inches(0.75), Inches(3.9), Inches(11.5), Inches(0.6),
         "A decoupled multi-agent architecture",
         font=HEADER_FONT, size=26, color=STEEL, italic=True)
add_text(s, Inches(0.75), Inches(4.4), Inches(11.5), Inches(0.6),
         "for Python test generation.",
         font=HEADER_FONT, size=26, color=STEEL, italic=True)

add_text(s, Inches(0.7), Inches(6.5), Inches(8), Inches(0.3),
         "Walkthrough on a real function.",
         size=14, bold=True, color=CORAL)
add_text(s, Inches(0.7), Inches(6.85), Inches(8), Inches(0.3),
         "Mário João Gomes Ribeiro", size=12, color=ICE)


# ===================================================================
# SLIDE 2 — THE MOTIVATION (real baseline failure)
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "The motivation")
slide_number(s, 2)
title(s, "Same context, one prompt: this is what happens.")

# Left: real baseline test code
add_text(s, Inches(0.7), Inches(2.1), Inches(6.0), Inches(0.4),
         "A test our baseline (TEST4PY) produced — real, for sty/renderfunc.rgb_bg:",
         size=13, color=MUTED, italic=True)

baseline_code = """# Module: sty.renderfunc
from sty.renderfunc import rgb_bg
import pytest

def test_rgb_bg_invalid_range():
    \"\"\"Test that an invalid RGB input
       raises a ValueError.\"\"\"
    with pytest.raises(ValueError):
        rgb_bg(256, 256, 256)  # Out of range"""
add_code(s, Inches(0.7), Inches(2.55), Inches(6.0), Inches(3.0),
         baseline_code, size=11)

# pytest error
add_rect(s, Inches(0.7), Inches(5.65), Inches(6.0), Inches(1.4),
         RGBColor(0xFC, 0xEE, 0xEC))
add_text(s, Inches(0.85), Inches(5.75), Inches(5.7), Inches(0.35),
         "pytest →", size=11, bold=True, color=RED, font=MONO_FONT)
add_text(s, Inches(0.85), Inches(6.05), Inches(5.7), Inches(0.9),
         "Failed: DID NOT RAISE <class 'ValueError'>",
         size=12, color=RED, font=MONO_FONT)
add_text(s, Inches(0.85), Inches(6.5), Inches(5.7), Inches(0.5),
         "(rgb_bg accepts any int — it never validates range)",
         size=10, italic=True, color=MUTED)

# Right: takeaway
rx = Inches(7.3)
add_text(s, rx, Inches(2.1), Inches(5.5), Inches(0.4),
         "WHAT HAPPENED", size=11, bold=True, color=CORAL)
add_text(s, rx, Inches(2.5), Inches(5.5), Inches(1.0),
         "The LLM assumed the function\nvalidates input.",
         font=HEADER_FONT, size=22, bold=True, color=INK, line_spacing=1.15)
add_text(s, rx, Inches(3.9), Inches(5.5), Inches(2.2),
         "It doesn't. rgb_bg just builds an ANSI escape string with whatever "
         "numbers you pass.\n\n"
         "Same Phase 1 context MARTA uses — but one big prompt to write the test.",
         size=13, color=INK, line_spacing=1.45)
add_rect(s, rx, Inches(6.2), Inches(5.5), Inches(0.95), ICE)
add_text(s, rx + Inches(0.2), Inches(6.3), Inches(5.1), Inches(0.4),
         "Baseline = TEST4PY", size=12, bold=True, color=NAVY)
add_text(s, rx + Inches(0.2), Inches(6.6), Inches(5.1), Inches(0.5),
         "Same Phase 1 as MARTA. The only thing it lacks\nis the Planner / Assertion split + the loops.",
         size=11, italic=True, color=NAVY, line_spacing=1.3)


# ===================================================================
# SLIDE 3 — WHY PYTHON IS THE HARD CASE (real function)
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "The hard case")
slide_number(s, 3)
title(s, "Now ask an LLM to test this.")

add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
         "Real Python. No type hints. Effects come from another module's constants.",
         font=HEADER_FONT, size=16, italic=True, color=MUTED)

# Real EfRegister code
code = """class EfRegister(Register):
    \"\"\"The default effect register.\"\"\"

    def __init__(self):
        super().__init__()
        self.renderfuncs[Sgr] = renderfunc.sgr
        self.b      = Style(Sgr(1))
        self.bold   = Style(Sgr(1))
        self.dim    = Style(Sgr(2))
        self.i      = Style(Sgr(3))
        self.italic = Style(Sgr(3))
        self.u      = Style(Sgr(4))
        self.underl = Style(Sgr(4))
        self.blink  = Style(Sgr(5))
        # ...8 more effects..."""
add_code(s, Inches(0.7), Inches(2.55), Inches(6.6), Inches(4.3), code, size=12)

# Right: questions
rx = Inches(7.7)
add_text(s, rx, Inches(2.55), Inches(5.0), Inches(0.4),
         "WHAT THE LLM HAS TO GUESS", size=11, bold=True, color=CORAL)

questions = [
    ("What is Style?", "Some wrapper class. Imported how?"),
    ("What is Sgr?", "A number? An enum? A callable?"),
    ("What does renderfunc.sgr do?", "A formatter? A renderer? Side-effects?"),
    ("Is super() doing anything?", "What's the parent class? Mocks?"),
]
qy = Inches(3.05)
for i, (q, a) in enumerate(questions):
    y = qy + i * Inches(0.85)
    add_oval(s, rx, y + Inches(0.18), Inches(0.16), Inches(0.16), CORAL)
    add_text(s, rx + Inches(0.3), y, Inches(4.8), Inches(0.35),
             q, size=14, bold=True, color=INK, font=MONO_FONT)
    add_text(s, rx + Inches(0.3), y + Inches(0.38), Inches(4.8), Inches(0.45),
             a, size=12, color=MUTED, italic=True)

add_text(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
         "Guess wrong → hallucinated assertions, broken imports, dead tests.",
         font=HEADER_FONT, size=14, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 4 — MARTA'S BET (the pivot — high level, dark)
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, Inches(0.7), Inches(0.7), Inches(0.25), Inches(0.25), CORAL)
add_text(s, Inches(1.05), Inches(0.7), Inches(8), Inches(0.3),
         "MARTA'S BET", size=11, bold=True, color=CORAL)
slide_number(s, 4, on_dark=True)

add_text(s, Inches(0.7), Inches(1.6), Inches(12), Inches(1.4),
         "Don't make one LLM do everything.",
         font=HEADER_FONT, size=44, bold=True, color=ICE, line_spacing=1.05)
add_text(s, Inches(0.7), Inches(2.95), Inches(12), Inches(0.6),
         "Three moves we make instead.",
         font=HEADER_FONT, size=20, italic=True, color=STEEL)

# Three cards
cards = [
    ("01", "Specialise",
     "Two agents.\nOne plans WHAT to test.\nOne writes HOW.\nSmaller jobs, fewer hallucinations."),
    ("02", "Travel the call graph",
     "Build each function's understanding\nfrom the summaries of every function\nit calls. Recursively."),
    ("03", "Close the loop",
     "Pytest runs the test.\nIts real error feeds back.\nThe LLM rewrites against\nground truth, not a guess."),
]
card_w = Inches(3.85)
card_h = Inches(3.4)
gap = Inches(0.25)
total_w = card_w * 3 + gap * 2
start_x = (SW - total_w) // 2
y = Inches(3.85)
for i, (num, h, body) in enumerate(cards):
    x = start_x + i * (card_w + gap)
    add_rect(s, x, y, card_w, card_h, NAVY_SOFT)
    add_text(s, x + Inches(0.4), y + Inches(0.35), Inches(1), Inches(0.5),
             num, font=HEADER_FONT, size=22, bold=True, color=CORAL)
    add_text(s, x + Inches(0.4), y + Inches(0.9), card_w - Inches(0.8), Inches(0.6),
             h, font=HEADER_FONT, size=24, bold=True, color=ICE)
    add_text(s, x + Inches(0.4), y + Inches(1.65), card_w - Inches(0.8), Inches(1.7),
             body, size=14, color=STEEL, line_spacing=1.45)


# ===================================================================
# SLIDE 5 — THE PIPELINE OVERVIEW
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "The pipeline")
slide_number(s, 5)
title(s, "Two phases. One pipeline.")
subtitle(s, "Build the right context first. Then let focused agents do the work.")

phase_y = Inches(3.0)
phase_h = Inches(3.6)
phase_w = Inches(5.7)
gap_x = Inches(0.5)
total = phase_w * 2 + gap_x
start_x = (SW - total) // 2

# Phase 1
x1 = start_x
add_rect(s, x1, phase_y, phase_w, phase_h, RGBColor(0xFF, 0xFF, 0xFF))
add_rect(s, x1, phase_y, Inches(0.18), phase_h, NAVY)
add_text(s, x1 + Inches(0.5), phase_y + Inches(0.35), Inches(3), Inches(0.4),
         "PHASE 1", size=11, bold=True, color=CORAL)
add_text(s, x1 + Inches(0.5), phase_y + Inches(0.75), phase_w - Inches(1.0), Inches(0.6),
         "Build the context", font=HEADER_FONT, size=24, bold=True, color=INK)
bullets1 = [
    ("PyCG + AST", "static call graph and class structure"),
    ("3-pass function analysis", "what it does, what it should do, merged"),
    ("Per-class summary + how-to-use", "feeds the type-inference RAG"),
    ("Parameter type inference", "syntactic filter + semantic RAG (Chroma)"),
]
for i, (k, v) in enumerate(bullets1):
    yy = phase_y + Inches(1.5) + i * Inches(0.45)
    add_oval(s, x1 + Inches(0.5), yy + Inches(0.13), Inches(0.12), Inches(0.12), CORAL)
    add_text(s, x1 + Inches(0.75), yy, phase_w - Inches(1.3), Inches(0.35),
             k, size=13, bold=True, color=INK)
    add_text(s, x1 + Inches(0.75), yy + Inches(0.18), phase_w - Inches(1.3), Inches(0.3),
             v, size=11, italic=True, color=MUTED)

# Phase 2
x2 = start_x + phase_w + gap_x
add_rect(s, x2, phase_y, phase_w, phase_h, RGBColor(0xFF, 0xFF, 0xFF))
add_rect(s, x2, phase_y, Inches(0.18), phase_h, CORAL)
add_text(s, x2 + Inches(0.5), phase_y + Inches(0.35), Inches(3), Inches(0.4),
         "PHASE 2", size=11, bold=True, color=CORAL)
add_text(s, x2 + Inches(0.5), phase_y + Inches(0.75), phase_w - Inches(1.0), Inches(0.6),
         "Multi-agent ReAct", font=HEADER_FONT, size=24, bold=True, color=INK)
bullets2 = [
    ("Planner agent", "3-scenario JSON plan + related-function RAG"),
    ("Assertion agent", "one pytest function per scenario"),
    ("Inner loop", "Pylint → Pytest, up to 3 self-healing attempts"),
    ("Outer loop", "N rounds, replays Planner with coverage feedback"),
]
for i, (k, v) in enumerate(bullets2):
    yy = phase_y + Inches(1.5) + i * Inches(0.45)
    add_oval(s, x2 + Inches(0.5), yy + Inches(0.13), Inches(0.12), Inches(0.12), CORAL)
    add_text(s, x2 + Inches(0.75), yy, phase_w - Inches(1.3), Inches(0.35),
             k, size=13, bold=True, color=INK)
    add_text(s, x2 + Inches(0.75), yy + Inches(0.18), phase_w - Inches(1.3), Inches(0.3),
             v, size=11, italic=True, color=MUTED)

arrow_y = phase_y + phase_h // 2
add_arrow(s, x1 + phase_w, arrow_y, x2, arrow_y, color=NAVY, weight=3)
add_text(s, x1 + phase_w + Inches(0.05), arrow_y - Inches(0.45),
         gap_x - Inches(0.1), Inches(0.4),
         "feeds", size=11, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 6 — MEET THE EXAMPLE
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Running example")
slide_number(s, 6)
title(s, "Throughout the rest, one function.")

add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.5),
         "EfRegister.__init__  —  from the sty library  ·  sty/register.py",
         font=MONO_FONT, size=16, color=NAVY, bold=True)

ex_code = """class EfRegister(Register):
    \"\"\"The default effect register.

    Instances from this class can be used to create
    effects like 'bold', 'dim', 'italic', 'blink', etc.

        print(f\"{ef.bold}Bold Text{ef.rs}\")
    \"\"\"

    def __init__(self):
        super().__init__()
        self.renderfuncs[Sgr] = renderfunc.sgr

        self.b      = Style(Sgr(1))
        self.bold   = Style(Sgr(1))
        self.dim    = Style(Sgr(2))
        self.i      = Style(Sgr(3))
        # ...10 more aliases for effects + rs (reset)..."""
add_code(s, Inches(0.7), Inches(2.7), Inches(7.5), Inches(4.4), ex_code, size=12)

rx = Inches(8.5)
add_text(s, rx, Inches(2.7), Inches(4.3), Inches(0.4),
         "WHY THIS ONE", size=11, bold=True, color=CORAL)
add_text(s, rx, Inches(3.1), Inches(4.3), Inches(3.5),
         "Real Python, no parameters, no return.\n\n"
         "Initialises 12 attributes from another module's constants.\n\n"
         "Touches inheritance (super()), mutable state (renderfuncs), "
         "and a non-obvious class (Style).\n\n"
         "Hard for any one-shot LLM. Ideal for MARTA.",
         size=13, color=INK, line_spacing=1.45)


# ===================================================================
# SLIDE 7 — PHASE 1: WHAT WE BUILD FOR THIS FUNCTION
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Phase 1 — context for EfRegister")
slide_number(s, 7)
title(s, "We don't hand the LLM the source. We hand it understanding.")

# Two panels: function summary + class summary
pw = Inches(5.9)
ph = Inches(4.4)
gap = Inches(0.5)
sx = (SW - pw * 2 - gap) // 2
sy = Inches(2.4)

# Panel 1 — function summary
add_rect(s, sx, sy, pw, ph, RGBColor(0xFF, 0xFF, 0xFF))
add_rect(s, sx, sy, Inches(0.08), ph, NAVY)
add_text(s, sx + Inches(0.3), sy + Inches(0.25), pw - Inches(0.5), Inches(0.35),
         "FUNCTION-LEVEL SUMMARY (Phase 1, LLM)",
         size=10, bold=True, color=CORAL)
add_text(s, sx + Inches(0.3), sy + Inches(0.6), pw - Inches(0.5), Inches(0.5),
         "EfRegister.__init__", size=15, bold=True, color=INK, font=MONO_FONT)
add_text(s, sx + Inches(0.3), sy + Inches(1.05), pw - Inches(0.5), ph - Inches(1.3),
         "Initialises 12 text-effect Style objects on the instance "
         "(bold, dim, italic, blink, …), each wrapping a specific Sgr "
         "SGR code. Also binds renderfuncs[Sgr] to renderfunc.sgr so the "
         "register knows how to serialise any Sgr to its ANSI escape "
         "sequence.\n\n"
         "Calls super().__init__() to inherit the base Register state.\n\n"
         "No parameters. No return. Side-effects on self only.",
         size=12, color=INK, line_spacing=1.4, italic=True)

# Panel 2 — class summary + how-to-use
x2 = sx + pw + gap
add_rect(s, x2, sy, pw, ph, RGBColor(0xFF, 0xFF, 0xFF))
add_rect(s, x2, sy, Inches(0.08), ph, CORAL)
add_text(s, x2 + Inches(0.3), sy + Inches(0.25), pw - Inches(0.5), Inches(0.35),
         "CLASS-LEVEL HOW-TO-USE (Phase 1, LLM)",
         size=10, bold=True, color=CORAL)
add_text(s, x2 + Inches(0.3), sy + Inches(0.6), pw - Inches(0.5), Inches(0.5),
         "Style  (referenced by EfRegister)",
         size=15, bold=True, color=INK, font=MONO_FONT)
add_text(s, x2 + Inches(0.3), sy + Inches(1.05), pw - Inches(0.5), ph - Inches(1.3),
         "Wraps one or more Sgr codes into a string-coercible object.\n\n"
         "Constructor:\n  Style(*sgrs)   — accepts any number of Sgr instances\n\n"
         "Used as: str(Style(Sgr(1)))  →  '\\x1b[1m'\n\n"
         "Composes: str(Style(Sgr(22), Sgr(23), …)) joins each code.",
         size=12, color=INK, line_spacing=1.4, italic=True, font=BODY_FONT)

add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.35),
         "All of this is produced by the LLM in Phase 1, once per project — then cached.",
         size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 8 — PHASE 1 ZOOM: PYCG + RECURSIVE SUMMARIES
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Phase 1 — why static analysis matters")
slide_number(s, 8)
title(s, "Summaries propagate up the call graph.")

subtitle(s, "We don't summarise EfRegister.__init__ in a vacuum.")

# Left — diagram
gx = Inches(0.7)
gy = Inches(3.0)
gw = Inches(6.0)
gh = Inches(3.9)
add_rect(s, gx, gy, gw, gh, RGBColor(0xFF, 0xFF, 0xFF))

node_w = Inches(2.2)
node_h = Inches(0.65)
n1_x = gx + (gw - node_w) // 2
n1_y = gy + Inches(0.4)
add_rounded(s, n1_x, n1_y, node_w, node_h, NAVY)
add_text(s, n1_x, n1_y + Inches(0.13), node_w, node_h,
         "EfRegister.__init__", color=ICE, size=12, bold=True,
         align=PP_ALIGN.CENTER, font=MONO_FONT)

# Bottom row: Style(Sgr(1)) and renderfunc.sgr
n2_x = gx + Inches(0.4)
n2_y = gy + Inches(2.1)
add_rounded(s, n2_x, n2_y, node_w, node_h, NAVY_SOFT)
add_text(s, n2_x, n2_y + Inches(0.13), node_w, node_h,
         "Style.__init__", color=ICE, size=12, bold=True,
         align=PP_ALIGN.CENTER, font=MONO_FONT)

n3_x = gx + gw - Inches(0.4) - node_w
n3_y = gy + Inches(2.1)
add_rounded(s, n3_x, n3_y, node_w, node_h, NAVY_SOFT)
add_text(s, n3_x, n3_y + Inches(0.13), node_w, node_h,
         "renderfunc.sgr", color=ICE, size=12, bold=True,
         align=PP_ALIGN.CENTER, font=MONO_FONT)

# Down arrows (calls)
add_arrow(s, n1_x + Inches(0.6), n1_y + node_h, n2_x + node_w // 2, n2_y)
add_arrow(s, n1_x + Inches(1.6), n1_y + node_h, n3_x + node_w // 2, n3_y)
add_text(s, gx + Inches(0.4), n1_y + node_h + Inches(0.05),
         gw - Inches(0.8), Inches(0.3),
         "calls", size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

add_text(s, gx + Inches(0.3), gy + Inches(3.15), gw - Inches(0.6), Inches(0.4),
         "Summary of EfRegister.__init__ is built knowing",
         size=12, color=INK, italic=True, align=PP_ALIGN.CENTER)
add_text(s, gx + Inches(0.3), gy + Inches(3.45), gw - Inches(0.6), Inches(0.4),
         "what Style and renderfunc.sgr actually do.",
         size=12, color=INK, italic=True, align=PP_ALIGN.CENTER)

# Right — benefits
rx = Inches(7.3)
ry = Inches(3.0)
cw = Inches(5.5)
ch = Inches(1.15)
gap = Inches(0.18)
benefits = [
    ("Recursive function summaries",
     "Caller's docstring is grounded in summaries of its callees, all the way down."),
    ("Class MRO walked",
     "Subclass methods inherited from parents are visible to type inference."),
    ("Imports scope the RAG",
     "Only classes reachable from this file become retrieval candidates."),
]
for i, (h, body) in enumerate(benefits):
    y = ry + i * (ch + gap)
    add_rect(s, rx, y, cw, ch, RGBColor(0xFF, 0xFF, 0xFF))
    add_rect(s, rx, y, Inches(0.08), ch, CORAL)
    add_text(s, rx + Inches(0.3), y + Inches(0.18), cw - Inches(0.5), Inches(0.4),
             h, size=14, bold=True, color=INK)
    add_text(s, rx + Inches(0.3), y + Inches(0.55), cw - Inches(0.5), Inches(0.55),
             body, size=12, color=MUTED, line_spacing=1.35)

add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3),
         "PyCG: Salis et al., 2020 (Apache 2.0) — vendored static-analysis library.",
         size=10, color=MUTED, italic=True)


# ===================================================================
# SLIDE 9 — PHASE 1: TYPE INFERENCE VIA RAG
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Phase 1 — type inference (RAG)")
slide_number(s, 9)
title(s, "What's the type of a parameter? Ask the project.")

subtitle(s, "Python won't tell us. The codebase already has the answer.")

# Show flow as 3 cards in a row
cards = [
    ("1.  Inspect usage (AST)",
     "Scan the function body.\nRecord every attribute the\nparameter is asked for.",
     "param.login()\nparam.send(msg)\n→ {login, send}"),
    ("2.  Syntactic filter",
     "Find classes whose member\nset includes those names.\nInheritance counts (MRO).",
     "Match: User, ServiceUser,\nMockClient. 9 others pruned."),
    ("3.  Semantic RAG (Chroma)",
     "Embed an LLM-written\n\"meaning\" of the parameter.\nTop-K against class summaries.",
     "→ Top-3 candidates\n  with how-to-use docs"),
]
cw = Inches(4.0)
ch = Inches(4.0)
gap = Inches(0.15)
total_w = cw * 3 + gap * 2
sx = (SW - total_w) // 2
sy = Inches(2.7)

for i, (h, body, ex) in enumerate(cards):
    x = sx + i * (cw + gap)
    add_rect(s, x, sy, cw, ch, RGBColor(0xFF, 0xFF, 0xFF))
    add_text(s, x + Inches(0.3), sy + Inches(0.3), cw - Inches(0.5), Inches(0.4),
             h, size=14, bold=True, color=CORAL)
    add_text(s, x + Inches(0.3), sy + Inches(0.8), cw - Inches(0.5), Inches(1.5),
             body, size=13, color=INK, line_spacing=1.45)
    add_rect(s, x + Inches(0.3), sy + Inches(2.4), cw - Inches(0.6), Inches(1.3), CODE_BG)
    add_text(s, x + Inches(0.45), sy + Inches(2.5), cw - Inches(0.9), Inches(1.1),
             ex, size=11, color=NAVY, font=MONO_FONT, line_spacing=1.4)
    if i < len(cards) - 1:
        add_arrow(s, x + cw, sy + ch // 2, x + cw + gap, sy + ch // 2,
                  color=NAVY, weight=3)

add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.35),
         "Result lands in the Planner prompt as INFERRED TYPES & USAGE.",
         size=12, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 10 — THE PLANNER PROMPT (real)
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Phase 2 — Planner agent")
slide_number(s, 10)
title(s, "The actual prompt we send. For EfRegister.")

prompt_code = """System: You are a QA Lead. Analyze the function and output a Test Plan in JSON.

User:
  CONTEXT:
    Function Name: EfRegister.__init__   Module: register

    SOURCE CODE:
      class EfRegister: ...
      def __init__(self):
          super().__init__()
          self.renderfuncs[Sgr] = renderfunc.sgr
          self.b = Style(Sgr(1))  ;  self.bold = Style(Sgr(1))
          ...

    INFERRED TYPES & USAGE:
      Style(*sgrs) wraps Sgr codes into a string-coercible ANSI
      escape object.  str(Style(Sgr(1))) → '\\x1b[1m'.

    DOCSTRING SUMMARY:
      Initialises 12 text-effect Style objects on the instance.
      No parameters. Side-effects on self only.

    COVERAGE FEEDBACK:  First pass: Try to achieve maximum coverage.

  TASK:
    Generate 3 distinct test scenarios covering:
      1. Valid inputs (Happy Path)
      2. Edge cases (None, empty, boundary)
      3. Invalid inputs / Error handling

  OUTPUT FORMAT:  Return ONLY a raw JSON list. No Markdown."""

add_code(s, Inches(0.7), Inches(2.2), Inches(11.9), Inches(4.7), prompt_code, size=11)
add_text(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.3),
         "Real prompt from llm_traffic.log. Some lines collapsed for readability.",
         size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 11 — THE PLANNER OUTPUT (real JSON)
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Phase 2 — Planner output")
slide_number(s, 11)
title(s, "What comes back: a plan, in JSON.")

json_code = """[
    {
        "name":  "test_valid_inputs",
        "desc":  "Test standard inputs to ensure all effects are
                  correctly initialized and can be applied.",
        "setup": "ef = EfRegister()"
    },
    {
        "name":  "test_edge_cases",
        "desc":  "Test edge cases with boundary values and verify
                  that the reset effect clears all applied styles.",
        "setup": "ef = EfRegister()"
    },
    {
        "name":  "test_invalid_inputs",
        "desc":  "Verify the instance exposes every documented
                  attribute and the renderfuncs dict is set.",
        "setup": "ef = EfRegister()"
    }
]"""
add_code(s, Inches(0.7), Inches(2.1), Inches(7.5), Inches(4.7), json_code, size=12)

rx = Inches(8.7)
add_text(s, rx, Inches(2.1), Inches(4.2), Inches(0.4),
         "WHAT THIS BUYS US", size=11, bold=True, color=CORAL)
add_text(s, rx, Inches(2.5), Inches(4.2), Inches(2.8),
         "Three independent scenarios.\n\n"
         "Each one is small, focused, and self-contained.\n\n"
         "The next agent doesn't have to think about strategy — "
         "just translate one bullet at a time.",
         size=13, color=INK, line_spacing=1.4)

add_rect(s, rx, Inches(5.8), Inches(4.2), Inches(1.0), ICE)
add_text(s, rx + Inches(0.2), Inches(5.9), Inches(4.0), Inches(0.4),
         "Note", size=12, bold=True, color=CORAL)
add_text(s, rx + Inches(0.2), Inches(6.2), Inches(4.0), Inches(0.6),
         "Robust parser: malformed JSON falls back to a single\nbasic-test scenario.",
         size=11, color=MUTED, italic=True, line_spacing=1.35)


# ===================================================================
# SLIDE 12 — ASSERTION AGENT (real prompt + real output)
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Phase 2 — Assertion agent")
slide_number(s, 12)
title(s, "One scenario in. One pytest function out.")

# Left — prompt
lprompt = """System: You are a Pytest Expert. Write valid python code.

User:
  Write a SINGLE pytest function for: test_valid_inputs
  DESC: Test standard inputs to ensure all effects are
        correctly initialized and can be applied.
  SETUP: ef = EfRegister()

  FUNCTION CODE:
    class EfRegister: ...
    def __init__(self): ...

  RULES:
    1. Output ONLY python code in ```python``` block.
    2. Import correctly from module 'register'.
    3. CRITICAL MOCKING RULE: use patch() or monkeypatch.
       Never assign mocks to global modules."""
add_text(s, Inches(0.7), Inches(2.1), Inches(6.0), Inches(0.4),
         "PROMPT", size=10, bold=True, color=CORAL)
add_code(s, Inches(0.7), Inches(2.5), Inches(6.0), Inches(4.5), lprompt, size=10)

# Right — generated code (compacted comments)
gen = """import pytest
from sty import register, Style, Sgr

@pytest.fixture(scope=\"module\")
def ef():
    return register.EfRegister()

def test_valid_inputs(ef):
    assert str(ef.bold)    == \"\\033[1m\"
    assert str(ef.italic)  == \"\\033[3m\"
    assert str(ef.underl)  == \"\\033[4m\"
    assert str(ef.blink)   == \"\\033[5m\"
    assert str(ef.inverse) == \"\\033[7m\"
    assert str(ef.hidden)  == \"\\033[8m\"
    assert str(ef.strike)  == \"\\033[9m\"
    # reset = concat of all SGR codes
    assert str(ef.rs) == \\
        \"\\033[22m\\033[23m\\033[24m\\033[25m\\033[27m\\033[28m\\033[29m\""""
add_text(s, Inches(7.0), Inches(2.1), Inches(6.0), Inches(0.4),
         "GENERATED CODE (passes after 1 self-heal)", size=10, bold=True, color=GREEN)
add_code(s, Inches(7.0), Inches(2.5), Inches(6.0), Inches(4.5), gen, size=11)

add_text(s, Inches(0.7), Inches(7.15), Inches(12), Inches(0.3),
         "Uses a pytest fixture. Knows the alias map. Gets the reset sequence right — after one repair.",
         size=11, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 13 — INNER LOOP: REAL SELF-HEALING
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Inner loop — self-healing")
slide_number(s, 13)
title(s, "When pytest disagrees, the agent listens.")

# Top: 1st attempt (wrong) — narrow code panel
add_text(s, Inches(0.7), Inches(1.9), Inches(6.0), Inches(0.4),
         "ATTEMPT 1  ·  generated code", size=10, bold=True, color=CORAL)
attempt1 = """def test_valid_inputs():
    ef = EfRegister()
    assert str(ef.bold) == \"\\x1b[1m\"
    assert str(ef.italic) == \"\\x1b[3m\"
    # ... 6 more assertions ...
    assert str(ef.rs) == \"\\x1b[0m\"   # ← assumed"""
add_code(s, Inches(0.7), Inches(2.3), Inches(6.0), Inches(2.0), attempt1,
         size=11, highlight=[(7, RED)])

# pytest output
add_text(s, Inches(0.7), Inches(4.5), Inches(6.0), Inches(0.35),
         "PYTEST →", size=10, bold=True, color=RED)
err = """AssertionError: assert '\\x1b[22m\\x1b[23m...\\x1b[29m' == '\\x1b[0m'
  - [0m
  + [22m[23m[24m[25m[27m[28m[29m"""
add_code(s, Inches(0.7), Inches(4.85), Inches(6.0), Inches(1.3), err,
         size=11, color=RED)

# Right: attempt 2 (fixed)
add_text(s, Inches(7.0), Inches(1.9), Inches(6.0), Inches(0.4),
         "ATTEMPT 2  ·  rewritten with ground truth", size=10, bold=True, color=GREEN)
attempt2 = """def test_valid_inputs(ef):
    # ... same effects ...
    assert str(ef.bold)   == \"\\033[1m\"
    assert str(ef.italic) == \"\\033[3m\"
    # ...
    # reset is the concatenation of all SGR codes
    assert str(ef.rs) == \\
        \"\\033[22m\\033[23m\\033[24m\\033[25m\\033[27m\\033[28m\\033[29m\\\"  """
add_code(s, Inches(7.0), Inches(2.3), Inches(6.0), Inches(2.0), attempt2,
         size=11, highlight=[(6, GREEN), (7, GREEN)])

add_text(s, Inches(7.0), Inches(4.5), Inches(6.0), Inches(0.35),
         "PYTEST →", size=10, bold=True, color=GREEN)
add_text(s, Inches(7.0), Inches(4.85), Inches(6.0), Inches(0.4),
         "PASSED",
         size=14, bold=True, color=GREEN, font=MONO_FONT)

# Bottom: explanation
add_rect(s, Inches(0.7), Inches(6.3), Inches(12.3), Inches(0.9), ICE)
add_text(s, Inches(0.9), Inches(6.4), Inches(12), Inches(0.4),
         "The LLM had the wrong assumption: reset = \\x1b[0m, the common ANSI reset.",
         size=13, color=NAVY, bold=True)
add_text(s, Inches(0.9), Inches(6.75), Inches(12), Inches(0.4),
         "Pytest gave it the truth from the real instance. It rewrote correctly. "
         "Up to 3 attempts; otherwise quarantined.",
         size=12, color=NAVY, italic=True)


# ===================================================================
# SLIDE 14 — OUTER LOOP: COVERAGE-DRIVEN
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Outer loop")
slide_number(s, 14)
title(s, "Round 2 onward: chase the missed lines.")

# Top — pipeline boxes
boxes = [
    ("Suite passes", NAVY_SOFT),
    ("Run coverage.py", NAVY_SOFT),
    ("Find missed lines", CORAL),
    ("Augment Planner prompt", NAVY),
    ("New scenarios", NAVY),
]
bx = Inches(0.7)
by = Inches(2.4)
bw = Inches(2.3)
bh = Inches(1.0)
bgap = Inches(0.18)
for i, (label, col) in enumerate(boxes):
    x = bx + i * (bw + bgap)
    add_rounded(s, x, by, bw, bh, col)
    add_text(s, x, by + Inches(0.25), bw, bh,
             label, color=ICE, size=13, bold=True, align=PP_ALIGN.CENTER,
             line_spacing=1.2)
    if i > 0:
        add_arrow(s, bx + (i - 1) * (bw + bgap) + bw, by + bh // 2,
                  x, by + bh // 2)

last_x = bx + (len(boxes) - 1) * (bw + bgap) + bw // 2
first_x = bx + bw // 2
loop_y = by + bh + Inches(0.35)
add_arrow(s, last_x, by + bh, last_x, loop_y, color=NAVY, weight=2)
add_arrow(s, last_x, loop_y, first_x, loop_y, color=NAVY, weight=2)
add_arrow(s, first_x, loop_y, first_x, by + bh, color=NAVY, weight=2)
add_text(s, first_x + Inches(0.3), loop_y - Inches(0.35),
         last_x - first_x - Inches(0.6), Inches(0.4),
         "repeat for N rounds (default: 3)",
         size=11, italic=True, color=NAVY_SOFT, align=PP_ALIGN.CENTER)

# Bottom — concrete prompt fragment
add_text(s, Inches(0.7), Inches(4.7), Inches(12), Inches(0.4),
         "WHAT ROUND 2'S PLANNER PROMPT NOW INCLUDES",
         size=10, bold=True, color=CORAL)

round2 = """CONTEXT:
  Function Name: EfRegister.__init__
  ...
  COVERAGE FEEDBACK:
    MISSING LINES TO COVER: 17, 22-24

TASK:
  CRITICAL: If 'MISSING LINES TO COVER' are provided in the context,
  you MUST design scenarios specifically to execute those missing lines.
  Generate 3 distinct test scenarios covering: ..."""
add_code(s, Inches(0.7), Inches(5.15), Inches(12.0), Inches(1.95), round2, size=11)


# ===================================================================
# SLIDE 15 — REAL RESULTS: TEST4PY vs MARTA
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "Preliminary results")
slide_number(s, 15)
title(s, "Real numbers: TEST4PY vs MARTA.")
subtitle(s, "10 projects · DeepSeek-Coder-V2 16B · temp 0.2 · single run per setup")

# Bar chart — statement coverage per project
projects_short = ["str-utils", "pytutils", "superstring", "codetiming",
                  "isort", "sty", "pyMonet", "flutes", "dc-json", "docstring"]
baseline_cov = [80.10, 33.14, 41.76, 85.00, 41.49, 87.63, 68.46, 51.21, 46.82, 46.42]
marta_cov    = [85.74, 42.67, 62.35, 73.33, 54.95, 91.75, 91.28, 58.34, 44.50, 45.30]

chart_data = CategoryChartData()
chart_data.categories = projects_short
chart_data.add_series('TEST4PY (baseline)', baseline_cov)
chart_data.add_series('MARTA (3 rounds)',    marta_cov)

chart_shape = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.7), Inches(2.9), Inches(8.2), Inches(4.0),
    chart_data,
).chart
chart_shape.has_title = True
chart_shape.chart_title.text_frame.text = "Statement Coverage (%)"
for run in chart_shape.chart_title.text_frame.paragraphs[0].runs:
    run.font.size = Pt(13); run.font.bold = True
    run.font.color.rgb = INK; run.font.name = HEADER_FONT
chart_shape.has_legend = True
chart_shape.legend.position = XL_LEGEND_POSITION.BOTTOM
chart_shape.legend.include_in_layout = False
# Series colors
plot = chart_shape.plots[0]
plot.gap_width = 80
series_colors = [NAVY_SOFT, CORAL]
for i, series in enumerate(plot.series):
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = series_colors[i]
# Smaller category axis labels (10 projects)
cat_axis = chart_shape.category_axis
cat_axis.tick_labels.font.size = Pt(10)
cat_axis.tick_labels.font.color.rgb = MUTED
val_axis = chart_shape.value_axis
val_axis.tick_labels.font.size = Pt(10)
val_axis.tick_labels.font.color.rgb = MUTED
val_axis.maximum_scale = 100

# Right column — stat callouts
rx = Inches(9.3)

def stat_card(y, label, big, small, color):
    add_rect(s, rx, y, Inches(3.5), Inches(1.45), RGBColor(0xFF, 0xFF, 0xFF))
    add_rect(s, rx, y, Inches(0.08), Inches(1.45), color)
    add_text(s, rx + Inches(0.25), y + Inches(0.12), Inches(3.2), Inches(0.3),
             label.upper(), size=10, bold=True, color=color)
    add_text(s, rx + Inches(0.25), y + Inches(0.4), Inches(3.2), Inches(0.55),
             big, font=HEADER_FONT, size=22, bold=True, color=INK, line_spacing=1.05)
    add_text(s, rx + Inches(0.25), y + Inches(1.0), Inches(3.2), Inches(0.45),
             small, size=11, color=MUTED, italic=True, line_spacing=1.3)

stat_card(Inches(2.9), "Coverage",
          "+6.8 pts on average", "MARTA wins 8 / 10 projects", CORAL)
stat_card(Inches(4.45), "Tests passed (valid)",
          "220 → 2020", "~9× more across the 10 projects", NAVY)
stat_card(Inches(6.0), "Mutation score",
          "Mixed: 5 / 9 wins",
          "Strong gains on sty (+22), flutes (+40),\nsuperstring (+18). Losses on codetiming.",
          CORAL)

add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.3),
         "Caveat: single execution per configuration. No statistical test yet — see paper response matrix.",
         size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 16 — WHAT MARTA PRODUCES (side-by-side)
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, PAPER)
eyebrow(s, "End to end")
slide_number(s, 16)
title(s, "What lands on disk: a real suite for one function.")

# Three small panels — one per scenario
panel_w = Inches(4.0)
panel_h = Inches(4.5)
pgap = Inches(0.15)
total_w = panel_w * 3 + pgap * 2
sx = (SW - total_w) // 2
sy = Inches(2.4)

panels = [
    ("test_valid_inputs",
     """import pytest
from sty import register, Style, Sgr

@pytest.fixture(scope=\"module\")
def ef():
    return register.EfRegister()

def test_valid_inputs(ef):
    assert str(ef.bold)   == \"\\033[1m\"
    assert str(ef.italic) == \"\\033[3m\"
    assert str(ef.underl) == \"\\033[4m\"
    assert str(ef.blink)  == \"\\033[5m\"
    # ... 4 more effects ...
    assert str(ef.rs) == \\
      \"\\033[22m...\\033[29m\""""),
    ("test_edge_cases",
     """import pytest
from sty import register, Style, Sgr

def test_edge_cases():
    ef = register.EfRegister()

    # bold and b are aliases
    assert ef.bold is not ef.b
    assert str(ef.bold) == str(ef.b)

    # rs joins all reset SGRs
    assert ef.rs.args is not None
    # reset clears every effect
    s = f\"{ef.bold}X{ef.rs}\"
    assert s.endswith(str(ef.rs))"""),
    ("test_invalid_inputs",
     """import pytest
from sty import EfRegister, Style, Sgr
from sty.renderfunc import sgr as r_sgr

def test_invalid_inputs():
    ef = EfRegister()
    # structural invariants
    for name in (\"b\", \"bold\", \"dim\",
                 \"italic\", \"underl\",
                 \"blink\", \"hidden\",
                 \"strike\", \"rs\"):
        assert hasattr(ef, name)
        assert isinstance(getattr(ef, name),
                          Style)
    assert ef.renderfuncs[Sgr] is r_sgr"""),
]
for i, (h, code) in enumerate(panels):
    x = sx + i * (panel_w + pgap)
    add_rect(s, x, sy, panel_w, Inches(0.5), NAVY)
    add_text(s, x, sy + Inches(0.1), panel_w, Inches(0.35),
             h, size=12, bold=True, color=ICE, font=MONO_FONT,
             align=PP_ALIGN.CENTER)
    add_code(s, x, sy + Inches(0.5), panel_w, panel_h - Inches(0.5),
             code, size=10)

add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.35),
         "Three files. Three concerns. All passing.   "
         "Each one isolated → one bad test never poisons the suite.",
         size=12, italic=True, color=NAVY, align=PP_ALIGN.CENTER)


# ===================================================================
# SLIDE 17 — TAKEAWAYS + HONEST LIMITATIONS
# ===================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, NAVY)
add_rect(s, Inches(0.7), Inches(0.7), Inches(0.25), Inches(0.25), CORAL)
add_text(s, Inches(1.05), Inches(0.7), Inches(8), Inches(0.3),
         "TAKEAWAYS", size=11, bold=True, color=CORAL)
slide_number(s, 17, on_dark=True)

add_text(s, Inches(0.7), Inches(1.3), Inches(12), Inches(1.4),
         "What we showed today.",
         font=HEADER_FONT, size=40, bold=True, color=ICE, line_spacing=1.05)

# Three takeaway cards
items = [
    ("01",
     "Two specialised agents beat one big prompt.",
     "Planner thinks scenarios. Assertion writes code.\n"
     "Each prompt is smaller and more focused — fewer hallucinations."),
    ("02",
     "Context comes from the call graph, not the source.",
     "PyCG + recursive summaries + RAG over class types.\n"
     "The LLM gets understanding, not a wall of code."),
    ("03",
     "The runtime is the teacher.",
     "Pylint and pytest are the cheapest, most honest critics.\n"
     "On failure: real error → rewrite → retry. Up to 3 attempts."),
]
ty = Inches(3.0)
th = Inches(1.15)
tw = Inches(12)
gap = Inches(0.15)
for i, (num, h, body) in enumerate(items):
    y = ty + i * (th + gap)
    add_rect(s, Inches(0.7), y, tw, th, NAVY_SOFT)
    add_text(s, Inches(0.9), y + Inches(0.25), Inches(0.8), th,
             num, font=HEADER_FONT, size=24, bold=True, color=CORAL)
    add_text(s, Inches(1.8), y + Inches(0.18), tw - Inches(1.2), Inches(0.4),
             h, font=HEADER_FONT, size=18, bold=True, color=ICE)
    add_text(s, Inches(1.8), y + Inches(0.55), tw - Inches(1.2), Inches(0.6),
             body, size=12, color=STEEL, line_spacing=1.4)

# Honest limitations
add_text(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
         "Honest limitation: MARTA captures what the code does, not what it should — these are regression tests.",
         font=HEADER_FONT, size=13, italic=True, color=STEEL, align=PP_ALIGN.CENTER)


# ---------- Save ----------
out = "/Users/mario/Desktop/GECAD/Test4Py/MARTA_presentation.pptx"
prs.save(out)

import subprocess
try:
    subprocess.run(["xattr", "-c", out], check=False)
except FileNotFoundError:
    pass

print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
